import streamlit as st
from pptx import Presentation
import fitz  # PyMuPDF
import os

def convert_pdf_to_ppt(uploaded_file, output_path="converted_presentation.pptx"):
    # Open the uploaded PDF
    pdf_document = fitz.open(stream=uploaded_file.read(), filetype="pdf")
    presentation = Presentation()
    temp_image_paths = []  # Store paths of temporary images

    for page_number in range(len(pdf_document)):
        page = pdf_document[page_number]

        # Extract text and images
        text = page.get_text()
        images = page.get_images(full=True)

        # Create a new slide
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])

        # Add text to slide
        if text:
            textbox = slide.shapes.add_textbox(left=100000, top=100000, width=5000000, height=5000000)
            textbox.text = text

        # Add images to slide
        for image_index, image in enumerate(images):
            xref = image[0]
            base_image = pdf_document.extract_image(xref)
            image_bytes = base_image["image"]
            image_path = f"temp_image_{page_number}_{image_index}.png"
            temp_image_paths.append(image_path)  # Track the temporary image

            # Save the image temporarily
            with open(image_path, "wb") as img_file:
                img_file.write(image_bytes)

            # Add the image to the slide
            slide.shapes.add_picture(image_path, left=50000, top=50000, width=3000000, height=2000000)

    # Save the presentation
    presentation.save(output_path)
    pdf_document.close()

    # Delete all temporary images
    for temp_image_path in temp_image_paths:
        if os.path.exists(temp_image_path):
            os.remove(temp_image_path)

    return output_path

# Streamlit App
st.title("Convert PDF to PowerPoint")

# Add unique keys to the components
uploaded_file = st.file_uploader("Upload a PDF file", type="pdf", key="pdf_to_ppt_uploader")

if uploaded_file:
    if st.button("Convert to PowerPoint", key="pdf_to_ppt_button"):
        with st.spinner("Converting..."):
            output_file = convert_pdf_to_ppt(uploaded_file)
        st.success("PDF converted to PowerPoint successfully!")
        with open(output_file, "rb") as file:
            st.download_button(
                label="Download PowerPoint",
                data=file,
                file_name="converted_presentation.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                key="pdf_to_ppt_download_button"
            )
        # Clean up the output file after download
        os.remove(output_file)
